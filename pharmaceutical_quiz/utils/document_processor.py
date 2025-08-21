"""
Document processing utilities for the Pharmaceutical Quiz Module.
Handles extraction of text content from various file formats.
"""

import os
import logging
from typing import List, Dict, Optional
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
import spacy
from config.settings import MAX_FILE_SIZE_MB, ERROR_MESSAGES


class DocumentProcessor:
    """Handles processing of uploaded documents to extract relevant context."""
    
    def __init__(self):
        self.supported_extensions = {'.pdf', '.docx', '.pptx', '.xlsx'}
        self.max_file_size = MAX_FILE_SIZE_MB * 1024 * 1024  # Convert to bytes
        self.nlp = None
        self._initialize_spacy()
    
    def _initialize_spacy(self):
        """Initialize spaCy NLP model for text processing."""
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            logging.warning("spaCy model 'en_core_web_sm' not found. Advanced text processing disabled.")
            self.nlp = None
    
    def process_files(self, file_paths: List[str]) -> str:
        """
        Process multiple files and extract text content.
        
        Args:
            file_paths: List of file paths to process
            
        Returns:
            Combined text content from all files
        """
        combined_content = []
        
        for file_path in file_paths:
            try:
                if not os.path.exists(file_path):
                    logging.warning(f"File not found: {file_path}")
                    continue
                
                # Check file size
                if os.path.getsize(file_path) > self.max_file_size:
                    logging.warning(f"File too large: {file_path}")
                    continue
                
                # Extract content based on file type
                content = self._extract_content(file_path)
                if content:
                    file_name = os.path.basename(file_path)
                    combined_content.append(f"=== Content from {file_name} ===\n{content}\n")
                    logging.info(f"Successfully processed: {file_name}")
                
            except Exception as e:
                logging.error(f"Error processing {file_path}: {e}")
                continue
        
        return "\n".join(combined_content)
    
    def _extract_content(self, file_path: str) -> Optional[str]:
        """Extract text content from a single file."""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self._extract_pdf_content(file_path)
        elif file_extension == '.docx':
            return self._extract_docx_content(file_path)
        elif file_extension == '.pptx':
            return self._extract_pptx_content(file_path)
        elif file_extension == '.xlsx':
            return self._extract_xlsx_content(file_path)
        else:
            logging.warning(f"Unsupported file type: {file_extension}")
            return None
    
    def _extract_pdf_content(self, file_path: str) -> Optional[str]:
        """Extract text from PDF file using PyMuPDF (fitz)."""
        content_parts = []
        
        try:
            # Use PyMuPDF for superior text extraction
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    content_parts.append(f"Page {page_num + 1}:\n{text}\n")
            
            doc.close()
            
            if content_parts:
                full_text = "\n".join(content_parts)
                # Apply spaCy processing if available
                return self._process_with_spacy(full_text) if self.nlp else full_text
            
            return None
            
        except Exception as e:
            logging.error(f"PyMuPDF failed for {file_path}: {e}")
            return None
    
    def _extract_docx_content(self, file_path: str) -> Optional[str]:
        """Extract text from Word document."""
        try:
            doc = Document(file_path)
            content_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_content = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_content.append(cell.text.strip())
                    if row_content:
                        table_content.append(" | ".join(row_content))
                
                if table_content:
                    content_parts.append("Table:\n" + "\n".join(table_content))
            
            return "\n\n".join(content_parts) if content_parts else None
            
        except Exception as e:
            logging.error(f"Error extracting from Word document {file_path}: {e}")
            return None
    
    def _extract_pptx_content(self, file_path: str) -> Optional[str]:
        """Extract text from PowerPoint presentation."""
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"Slide {slide_num}:"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if len(slide_content) > 1:  # More than just the slide number
                    content_parts.append("\n".join(slide_content))
            
            return "\n\n".join(content_parts) if content_parts else None
            
        except Exception as e:
            logging.error(f"Error extracting from PowerPoint {file_path}: {e}")
            return None
    
    def _extract_xlsx_content(self, file_path: str) -> Optional[str]:
        """Extract text from Excel file."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # Skip empty sheets
                    if df.empty:
                        continue
                    
                    sheet_content = [f"Sheet: {sheet_name}"]
                    
                    # Add column headers
                    headers = [str(col) for col in df.columns if str(col) != 'nan']
                    if headers:
                        sheet_content.append("Columns: " + " | ".join(headers))
                    
                    # Add sample data (first few rows)
                    sample_rows = min(10, len(df))  # Limit to first 10 rows
                    for idx, row in df.head(sample_rows).iterrows():
                        row_data = []
                        for col in df.columns:
                            value = row[col]
                            if pd.notna(value) and str(value).strip():
                                row_data.append(f"{col}: {value}")
                        
                        if row_data:
                            sheet_content.append(" | ".join(row_data))
                    
                    if len(sheet_content) > 1:
                        content_parts.append("\n".join(sheet_content))
                
                except Exception as e:
                    logging.warning(f"Error processing sheet {sheet_name}: {e}")
                    continue
            
            return "\n\n".join(content_parts) if content_parts else None
            
        except Exception as e:
            logging.error(f"Error extracting from Excel file {file_path}: {e}")
            return None
    
    def get_file_info(self, file_path: str) -> Dict:
        """Get basic information about a file."""
        try:
            stat = os.stat(file_path)
            return {
                'name': os.path.basename(file_path),
                'size': stat.st_size,
                'size_mb': round(stat.st_size / (1024 * 1024), 2),
                'extension': os.path.splitext(file_path)[1].lower(),
                'supported': os.path.splitext(file_path)[1].lower() in self.supported_extensions
            }
        except Exception as e:
            logging.error(f"Error getting file info for {file_path}: {e}")
            return {'error': str(e)}
    
    def validate_files(self, file_paths: List[str]) -> Dict:
        """Validate a list of files for processing."""
        results = {
            'valid_files': [],
            'invalid_files': [],
            'errors': []
        }
        
        for file_path in file_paths:
            try:
                if not os.path.exists(file_path):
                    results['invalid_files'].append(file_path)
                    results['errors'].append(f"File not found: {os.path.basename(file_path)}")
                    continue
                
                file_info = self.get_file_info(file_path)
                
                if 'error' in file_info:
                    results['invalid_files'].append(file_path)
                    results['errors'].append(f"Error accessing {os.path.basename(file_path)}")
                    continue
                
                if not file_info['supported']:
                    results['invalid_files'].append(file_path)
                    results['errors'].append(f"Unsupported file type: {os.path.basename(file_path)}")
                    continue
                
                if file_info['size'] > self.max_file_size:
                    results['invalid_files'].append(file_path)
                    results['errors'].append(f"File too large: {os.path.basename(file_path)} ({file_info['size_mb']}MB)")
                    continue
                
                results['valid_files'].append(file_path)
                
            except Exception as e:
                results['invalid_files'].append(file_path)
                results['errors'].append(f"Error validating {os.path.basename(file_path)}: {str(e)}")
        
        return results
    
    def _process_with_spacy(self, text: str) -> str:
        """Process text with spaCy for entity extraction and enhancement."""
        if not self.nlp or not text:
            return text
        
        try:
            doc = self.nlp(text)
            
            # Extract entities
            entities = []
            for ent in doc.ents:
                if ent.label_ in ['ORG', 'PRODUCT', 'PERSON', 'GPE', 'MONEY', 'PERCENT']:
                    entities.append(f"{ent.text} ({ent.label_})")
            
            # Add entity information to the text
            if entities:
                entity_summary = "\n\nKey Entities Identified:\n" + "\n".join(set(entities))
                return text + entity_summary
            
            return text
            
        except Exception as e:
            logging.warning(f"spaCy processing failed: {e}")
            return text
    
    def extract_pharmaceutical_entities(self, text: str) -> Dict:
        """Extract pharmaceutical-specific entities from text."""
        if not self.nlp:
            return {}
        
        try:
            doc = self.nlp(text)
            
            entities = {
                'organizations': [],
                'products': [],
                'locations': [],
                'monetary_values': [],
                'percentages': []
            }
            
            for ent in doc.ents:
                if ent.label_ == 'ORG':
                    entities['organizations'].append(ent.text)
                elif ent.label_ in ['PRODUCT', 'SUBSTANCE']:
                    entities['products'].append(ent.text)
                elif ent.label_ == 'GPE':
                    entities['locations'].append(ent.text)
                elif ent.label_ == 'MONEY':
                    entities['monetary_values'].append(ent.text)
                elif ent.label_ == 'PERCENT':
                    entities['percentages'].append(ent.text)
            
            # Remove duplicates
            for key in entities:
                entities[key] = list(set(entities[key]))
            
            return entities
            
        except Exception as e:
            logging.warning(f"Entity extraction failed: {e}")
            return {}
